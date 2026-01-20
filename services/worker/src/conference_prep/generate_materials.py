"""
Conference Material Generation Module

Generates conference presentation materials:
- Poster PDF using reportlab (single page with structured sections)
- Slides PPTX using python-pptx (title, methods, results, conclusion slides)

Supports blinding mode to strip author/institution information for blind review.

Output directory: /data/artifacts/conference/<run_id>/
"""

from __future__ import annotations

import hashlib
import os
from dataclasses import dataclass, field
from datetime import datetime
from enum import Enum
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

# reportlab imports for PDF generation
try:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.pdfgen import canvas
    from reportlab.platypus import (
        Frame,
        PageTemplate,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

# python-pptx imports for slides generation
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# ============ Constants ============

DEFAULT_ARTIFACT_BASE = Path("/data/artifacts/conference")
POSTER_WIDTH_INCHES = 48
POSTER_HEIGHT_INCHES = 36
SLIDE_WIDTH_INCHES = 13.333  # 16:9 aspect ratio
SLIDE_HEIGHT_INCHES = 7.5


class MaterialType(str, Enum):
    """Types of conference materials that can be generated."""
    POSTER_PDF = "poster_pdf"
    SLIDES_PPTX = "slides_pptx"


# ============ Data Classes ============

@dataclass
class PosterContent:
    """Content for poster generation."""
    title: str
    authors: List[str] = field(default_factory=list)
    institutions: List[str] = field(default_factory=list)
    background: str = ""
    methods: str = ""
    results: str = ""
    conclusion: str = ""
    references: List[str] = field(default_factory=list)
    acknowledgments: str = ""
    contact_email: str = ""
    qr_code_url: Optional[str] = None

    def get_blinded(self) -> "PosterContent":
        """Return a copy with author/institution info stripped."""
        return PosterContent(
            title=self.title,
            authors=["[Authors Blinded for Review]"],
            institutions=["[Institutions Blinded for Review]"],
            background=self.background,
            methods=self.methods,
            results=self.results,
            conclusion=self.conclusion,
            references=self.references,
            acknowledgments="[Acknowledgments Blinded for Review]",
            contact_email="[Contact Blinded for Review]",
            qr_code_url=None,  # Remove QR code in blinded mode
        )


@dataclass
class SlideContent:
    """Content for slides generation."""
    title: str
    subtitle: Optional[str] = None
    authors: List[str] = field(default_factory=list)
    institutions: List[str] = field(default_factory=list)
    background: str = ""
    objectives: List[str] = field(default_factory=list)
    methods: str = ""
    methods_bullets: List[str] = field(default_factory=list)
    results: str = ""
    results_bullets: List[str] = field(default_factory=list)
    conclusion: str = ""
    conclusion_bullets: List[str] = field(default_factory=list)
    references: List[str] = field(default_factory=list)
    acknowledgments: str = ""
    contact_info: str = ""

    def get_blinded(self) -> "SlideContent":
        """Return a copy with author/institution info stripped."""
        return SlideContent(
            title=self.title,
            subtitle=self.subtitle,
            authors=["[Authors Blinded for Review]"],
            institutions=["[Institutions Blinded for Review]"],
            background=self.background,
            objectives=self.objectives,
            methods=self.methods,
            methods_bullets=self.methods_bullets,
            results=self.results,
            results_bullets=self.results_bullets,
            conclusion=self.conclusion,
            conclusion_bullets=self.conclusion_bullets,
            references=self.references,
            acknowledgments="[Acknowledgments Blinded for Review]",
            contact_info="[Contact Blinded for Review]",
        )


@dataclass
class MaterialGenerationInput:
    """Input for material generation."""
    run_id: str
    material_type: MaterialType
    poster_content: Optional[PosterContent] = None
    slide_content: Optional[SlideContent] = None
    blinded: bool = False
    poster_size: Tuple[float, float] = (POSTER_WIDTH_INCHES, POSTER_HEIGHT_INCHES)
    output_dir: Optional[Path] = None

    def get_output_dir(self) -> Path:
        """Get output directory, creating if needed."""
        if self.output_dir:
            return self.output_dir
        return DEFAULT_ARTIFACT_BASE / self.run_id


@dataclass
class MaterialGenerationResult:
    """Result of material generation."""
    status: str  # "success", "error"
    material_type: MaterialType
    output_path: Optional[Path] = None
    file_size_bytes: int = 0
    sha256_hash: str = ""
    blinded: bool = False
    error_message: Optional[str] = None
    generation_timestamp: str = ""
    tool_version: str = ""

    def to_dict(self) -> dict:
        """Convert to dictionary for JSON serialization."""
        return {
            "status": self.status,
            "material_type": self.material_type.value,
            "output_path": str(self.output_path) if self.output_path else None,
            "file_size_bytes": self.file_size_bytes,
            "sha256_hash": self.sha256_hash,
            "blinded": self.blinded,
            "error_message": self.error_message,
            "generation_timestamp": self.generation_timestamp,
            "tool_version": self.tool_version,
        }


# ============ PDF Poster Generation ============

def _get_poster_styles() -> Dict[str, ParagraphStyle]:
    """Get paragraph styles for poster generation."""
    if not REPORTLAB_AVAILABLE:
        return {}

    styles = getSampleStyleSheet()

    return {
        "title": ParagraphStyle(
            "PosterTitle",
            parent=styles["Heading1"],
            fontSize=36,
            leading=42,
            alignment=TA_CENTER,
            spaceAfter=12,
            textColor=colors.HexColor("#1a365d"),
        ),
        "authors": ParagraphStyle(
            "PosterAuthors",
            parent=styles["Normal"],
            fontSize=18,
            leading=22,
            alignment=TA_CENTER,
            spaceAfter=6,
            textColor=colors.HexColor("#2d3748"),
        ),
        "institutions": ParagraphStyle(
            "PosterInstitutions",
            parent=styles["Normal"],
            fontSize=14,
            leading=18,
            alignment=TA_CENTER,
            spaceAfter=20,
            textColor=colors.HexColor("#4a5568"),
            fontName="Helvetica-Oblique",
        ),
        "section_header": ParagraphStyle(
            "SectionHeader",
            parent=styles["Heading2"],
            fontSize=20,
            leading=24,
            spaceAfter=8,
            spaceBefore=16,
            textColor=colors.HexColor("#2b6cb0"),
            borderPadding=4,
            backColor=colors.HexColor("#ebf8ff"),
        ),
        "body": ParagraphStyle(
            "PosterBody",
            parent=styles["Normal"],
            fontSize=12,
            leading=16,
            alignment=TA_JUSTIFY,
            spaceAfter=8,
        ),
        "reference": ParagraphStyle(
            "PosterReference",
            parent=styles["Normal"],
            fontSize=9,
            leading=12,
            leftIndent=20,
            firstLineIndent=-20,
        ),
    }


def generate_poster_pdf(
    content: PosterContent,
    output_path: Path,
    poster_size: Tuple[float, float] = (POSTER_WIDTH_INCHES, POSTER_HEIGHT_INCHES),
    blinded: bool = False,
) -> MaterialGenerationResult:
    """
    Generate a poster PDF using reportlab.

    Args:
        content: Poster content (title, authors, sections, etc.)
        output_path: Path for output PDF
        poster_size: Poster dimensions in inches (width, height)
        blinded: If True, strip author/institution information

    Returns:
        MaterialGenerationResult with status and file info
    """
    if not REPORTLAB_AVAILABLE:
        return MaterialGenerationResult(
            status="error",
            material_type=MaterialType.POSTER_PDF,
            error_message="reportlab not available. Install with: pip install reportlab",
            generation_timestamp=datetime.utcnow().isoformat() + "Z",
            tool_version="reportlab:unavailable",
        )

    try:
        # Apply blinding if requested
        if blinded:
            content = content.get_blinded()

        # Ensure output directory exists
        output_path.parent.mkdir(parents=True, exist_ok=True)

        # Calculate page size in points (1 inch = 72 points)
        page_width = poster_size[0] * inch
        page_height = poster_size[1] * inch

        # Create PDF document
        doc = SimpleDocTemplate(
            str(output_path),
            pagesize=(page_width, page_height),
            leftMargin=0.75 * inch,
            rightMargin=0.75 * inch,
            topMargin=0.75 * inch,
            bottomMargin=0.75 * inch,
        )

        styles = _get_poster_styles()
        story = []

        # Title
        story.append(Paragraph(content.title, styles["title"]))
        story.append(Spacer(1, 12))

        # Authors
        if content.authors:
            authors_text = ", ".join(content.authors)
            story.append(Paragraph(authors_text, styles["authors"]))

        # Institutions
        if content.institutions:
            institutions_text = "; ".join(content.institutions)
            story.append(Paragraph(institutions_text, styles["institutions"]))

        story.append(Spacer(1, 24))

        # Create two-column layout data
        section_data = []

        # Background section
        if content.background:
            section_data.append(("BACKGROUND", content.background))

        # Methods section
        if content.methods:
            section_data.append(("METHODS", content.methods))

        # Results section
        if content.results:
            section_data.append(("RESULTS", content.results))

        # Conclusion section
        if content.conclusion:
            section_data.append(("CONCLUSIONS", content.conclusion))

        # Build sections
        for section_title, section_text in section_data:
            story.append(Paragraph(section_title, styles["section_header"]))
            story.append(Paragraph(section_text, styles["body"]))
            story.append(Spacer(1, 12))

        # References
        if content.references:
            story.append(Paragraph("REFERENCES", styles["section_header"]))
            for i, ref in enumerate(content.references, 1):
                story.append(Paragraph(f"[{i}] {ref}", styles["reference"]))

        # Acknowledgments
        if content.acknowledgments:
            story.append(Spacer(1, 16))
            story.append(Paragraph("ACKNOWLEDGMENTS", styles["section_header"]))
            story.append(Paragraph(content.acknowledgments, styles["body"]))

        # Build PDF
        doc.build(story)

        # Calculate file info
        file_size = output_path.stat().st_size
        sha256_hash = _compute_file_hash(output_path)

        return MaterialGenerationResult(
            status="success",
            material_type=MaterialType.POSTER_PDF,
            output_path=output_path,
            file_size_bytes=file_size,
            sha256_hash=sha256_hash,
            blinded=blinded,
            generation_timestamp=datetime.utcnow().isoformat() + "Z",
            tool_version=f"reportlab:4.2.0",
        )

    except Exception as e:
        return MaterialGenerationResult(
            status="error",
            material_type=MaterialType.POSTER_PDF,
            error_message=str(e),
            generation_timestamp=datetime.utcnow().isoformat() + "Z",
            tool_version="reportlab:4.2.0",
        )


# ============ PPTX Slides Generation ============

def _add_title_slide(prs: "Presentation", content: SlideContent) -> None:
    """Add title slide to presentation."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(12.33), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = content.title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0x1a, 0x36, 0x5d)
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle if present
    if content.subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.0), Inches(12.33), Inches(0.75)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = content.subtitle
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = RGBColor(0x4a, 0x55, 0x68)
        subtitle_para.alignment = PP_ALIGN.CENTER

    # Authors
    if content.authors:
        authors_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.5), Inches(12.33), Inches(0.75)
        )
        authors_frame = authors_box.text_frame
        authors_para = authors_frame.paragraphs[0]
        authors_para.text = ", ".join(content.authors)
        authors_para.font.size = Pt(20)
        authors_para.font.color.rgb = RGBColor(0x2d, 0x37, 0x48)
        authors_para.alignment = PP_ALIGN.CENTER

    # Institutions
    if content.institutions:
        inst_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.25), Inches(12.33), Inches(0.75)
        )
        inst_frame = inst_box.text_frame
        inst_para = inst_frame.paragraphs[0]
        inst_para.text = "; ".join(content.institutions)
        inst_para.font.size = Pt(16)
        inst_para.font.italic = True
        inst_para.font.color.rgb = RGBColor(0x71, 0x80, 0x96)
        inst_para.alignment = PP_ALIGN.CENTER


def _add_content_slide(
    prs: "Presentation",
    title: str,
    body_text: str = "",
    bullets: List[str] = None,
) -> None:
    """Add a content slide with title and body/bullets."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Slide title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(12.33), Inches(1.0)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0x2b, 0x6c, 0xb0)

    # Content area
    content_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(1.75), Inches(11.83), Inches(5.0)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    if body_text:
        para = content_frame.paragraphs[0]
        para.text = body_text
        para.font.size = Pt(20)
        para.font.color.rgb = RGBColor(0x2d, 0x37, 0x48)

    if bullets:
        for i, bullet_text in enumerate(bullets):
            if i == 0 and not body_text:
                para = content_frame.paragraphs[0]
            else:
                para = content_frame.add_paragraph()
            para.text = bullet_text
            para.font.size = Pt(20)
            para.font.color.rgb = RGBColor(0x2d, 0x37, 0x48)
            para.level = 0


def generate_slides_pptx(
    content: SlideContent,
    output_path: Path,
    blinded: bool = False,
) -> MaterialGenerationResult:
    """
    Generate presentation slides using python-pptx.

    Creates slides for: Title, Background/Objectives, Methods, Results, Conclusions

    Args:
        content: Slide content (title, sections, bullets, etc.)
        output_path: Path for output PPTX
        blinded: If True, strip author/institution information

    Returns:
        MaterialGenerationResult with status and file info
    """
    if not PPTX_AVAILABLE:
        return MaterialGenerationResult(
            status="error",
            material_type=MaterialType.SLIDES_PPTX,
            error_message="python-pptx not available. Install with: pip install python-pptx",
            generation_timestamp=datetime.utcnow().isoformat() + "Z",
            tool_version="python-pptx:unavailable",
        )

    try:
        # Apply blinding if requested
        if blinded:
            content = content.get_blinded()

        # Ensure output directory exists
        output_path.parent.mkdir(parents=True, exist_ok=True)

        # Create presentation with 16:9 aspect ratio
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
        prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)

        # Slide 1: Title slide
        _add_title_slide(prs, content)

        # Slide 2: Background/Objectives
        if content.background or content.objectives:
            _add_content_slide(
                prs,
                title="Background & Objectives",
                body_text=content.background if not content.objectives else "",
                bullets=content.objectives if content.objectives else None,
            )

        # Slide 3: Methods
        if content.methods or content.methods_bullets:
            _add_content_slide(
                prs,
                title="Methods",
                body_text=content.methods if not content.methods_bullets else "",
                bullets=content.methods_bullets if content.methods_bullets else None,
            )

        # Slide 4: Results
        if content.results or content.results_bullets:
            _add_content_slide(
                prs,
                title="Results",
                body_text=content.results if not content.results_bullets else "",
                bullets=content.results_bullets if content.results_bullets else None,
            )

        # Slide 5: Conclusions
        if content.conclusion or content.conclusion_bullets:
            _add_content_slide(
                prs,
                title="Conclusions",
                body_text=content.conclusion if not content.conclusion_bullets else "",
                bullets=content.conclusion_bullets if content.conclusion_bullets else None,
            )

        # Slide 6: Acknowledgments (optional, only if not blinded or has content)
        if content.acknowledgments:
            _add_content_slide(
                prs,
                title="Acknowledgments",
                body_text=content.acknowledgments,
            )

        # Slide 7: References (optional)
        if content.references:
            _add_content_slide(
                prs,
                title="References",
                bullets=content.references[:8],  # Limit to 8 references per slide
            )

        # Save presentation
        prs.save(str(output_path))

        # Calculate file info
        file_size = output_path.stat().st_size
        sha256_hash = _compute_file_hash(output_path)

        return MaterialGenerationResult(
            status="success",
            material_type=MaterialType.SLIDES_PPTX,
            output_path=output_path,
            file_size_bytes=file_size,
            sha256_hash=sha256_hash,
            blinded=blinded,
            generation_timestamp=datetime.utcnow().isoformat() + "Z",
            tool_version="python-pptx:1.0.2",
        )

    except Exception as e:
        return MaterialGenerationResult(
            status="error",
            material_type=MaterialType.SLIDES_PPTX,
            error_message=str(e),
            generation_timestamp=datetime.utcnow().isoformat() + "Z",
            tool_version="python-pptx:1.0.2",
        )


# ============ Utility Functions ============

def _compute_file_hash(file_path: Path) -> str:
    """Compute SHA256 hash of a file."""
    sha256_hash = hashlib.sha256()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            sha256_hash.update(chunk)
    return sha256_hash.hexdigest()


def generate_material(
    input_params: MaterialGenerationInput,
) -> MaterialGenerationResult:
    """
    Generate conference material based on input parameters.

    Main entry point for material generation.

    Args:
        input_params: Generation parameters including content and options

    Returns:
        MaterialGenerationResult with generated file info
    """
    output_dir = input_params.get_output_dir()
    output_dir.mkdir(parents=True, exist_ok=True)

    if input_params.material_type == MaterialType.POSTER_PDF:
        if not input_params.poster_content:
            return MaterialGenerationResult(
                status="error",
                material_type=MaterialType.POSTER_PDF,
                error_message="poster_content required for POSTER_PDF generation",
                generation_timestamp=datetime.utcnow().isoformat() + "Z",
                tool_version="unknown",
            )

        output_path = output_dir / f"poster_{input_params.run_id}.pdf"
        return generate_poster_pdf(
            content=input_params.poster_content,
            output_path=output_path,
            poster_size=input_params.poster_size,
            blinded=input_params.blinded,
        )

    elif input_params.material_type == MaterialType.SLIDES_PPTX:
        if not input_params.slide_content:
            return MaterialGenerationResult(
                status="error",
                material_type=MaterialType.SLIDES_PPTX,
                error_message="slide_content required for SLIDES_PPTX generation",
                generation_timestamp=datetime.utcnow().isoformat() + "Z",
                tool_version="unknown",
            )

        output_path = output_dir / f"slides_{input_params.run_id}.pptx"
        return generate_slides_pptx(
            content=input_params.slide_content,
            output_path=output_path,
            blinded=input_params.blinded,
        )

    else:
        return MaterialGenerationResult(
            status="error",
            material_type=input_params.material_type,
            error_message=f"Unknown material type: {input_params.material_type}",
            generation_timestamp=datetime.utcnow().isoformat() + "Z",
            tool_version="unknown",
        )


def check_dependencies() -> Dict[str, bool]:
    """Check availability of generation dependencies."""
    return {
        "reportlab": REPORTLAB_AVAILABLE,
        "python-pptx": PPTX_AVAILABLE,
    }


# ============ Demo/Fixture Content ============

def get_demo_poster_content() -> PosterContent:
    """Get demo poster content for testing."""
    return PosterContent(
        title="Machine Learning Prediction of Thyroid Nodule Malignancy Using Multimodal Clinical Data",
        authors=["Jane Smith, MD", "John Doe, PhD", "Emily Chen, MS"],
        institutions=[
            "Department of Surgery, University Hospital",
            "Department of Computer Science, Tech University",
        ],
        background=(
            "Thyroid nodules are common, affecting up to 65% of the population. "
            "While most are benign, 5-15% are malignant. Current diagnostic approaches "
            "rely on ultrasound characteristics and fine-needle aspiration cytology, "
            "but indeterminate results occur in 15-30% of cases. Machine learning "
            "approaches may improve diagnostic accuracy by integrating multiple data sources."
        ),
        methods=(
            "We developed a gradient boosting classifier using data from 2,847 patients "
            "with thyroid nodules evaluated at our institution (2018-2024). Features included "
            "ultrasound characteristics (TI-RADS scores, size, composition), clinical factors "
            "(age, sex, family history), and laboratory values (TSH, thyroglobulin). "
            "The model was validated using 5-fold cross-validation with AUC-ROC as the "
            "primary metric."
        ),
        results=(
            "The final model achieved AUC-ROC of 0.91 (95% CI: 0.88-0.94) on the validation set, "
            "compared to 0.82 for TI-RADS alone (p<0.001). At 95% sensitivity threshold, "
            "specificity was 68% (vs. 42% for TI-RADS). The model reduced unnecessary biopsies "
            "by 35% while maintaining 95% sensitivity for malignancy detection. Key predictive "
            "features included TI-RADS score, nodule size, and patient age."
        ),
        conclusion=(
            "Our machine learning model significantly improves thyroid nodule risk stratification "
            "compared to TI-RADS alone. Integration of this tool into clinical workflows could "
            "reduce unnecessary biopsies while maintaining high sensitivity for malignancy detection. "
            "External validation and prospective studies are needed."
        ),
        references=[
            "Haugen BR, et al. 2015 American Thyroid Association Guidelines. Thyroid. 2016;26(1):1-133.",
            "Tessler FN, et al. ACR Thyroid Imaging, Reporting and Data System (TI-RADS). JACR. 2017;14(5):587-595.",
            "Durante C, et al. The Diagnosis and Management of Thyroid Nodules. JAMA. 2018;319(9):914-924.",
        ],
        acknowledgments="Supported by NIH Grant R01-CA123456 and institutional research funds.",
        contact_email="jane.smith@university.edu",
    )


def get_demo_slide_content() -> SlideContent:
    """Get demo slide content for testing."""
    return SlideContent(
        title="Machine Learning Prediction of Thyroid Nodule Malignancy",
        subtitle="Using Multimodal Clinical Data",
        authors=["Jane Smith, MD", "John Doe, PhD", "Emily Chen, MS"],
        institutions=[
            "Department of Surgery, University Hospital",
            "Department of Computer Science, Tech University",
        ],
        background=(
            "Thyroid nodules are common, with 5-15% being malignant. "
            "Current diagnostic approaches have limitations with 15-30% indeterminate results."
        ),
        objectives=[
            "Develop ML model for thyroid nodule risk stratification",
            "Compare performance to TI-RADS classification",
            "Evaluate potential to reduce unnecessary biopsies",
        ],
        methods_bullets=[
            "Retrospective cohort: 2,847 patients (2018-2024)",
            "Features: ultrasound (TI-RADS), clinical, laboratory",
            "Gradient boosting classifier with 5-fold CV",
            "Primary outcome: AUC-ROC for malignancy prediction",
        ],
        results_bullets=[
            "AUC-ROC: 0.91 (95% CI: 0.88-0.94) vs 0.82 TI-RADS (p<0.001)",
            "At 95% sensitivity: 68% specificity (vs 42% TI-RADS)",
            "35% reduction in unnecessary biopsies",
            "Top features: TI-RADS score, nodule size, age",
        ],
        conclusion_bullets=[
            "ML model significantly improves risk stratification",
            "Potential to reduce unnecessary biopsies by 35%",
            "External validation needed before clinical implementation",
        ],
        references=[
            "Haugen BR, et al. Thyroid. 2016;26(1):1-133.",
            "Tessler FN, et al. JACR. 2017;14(5):587-595.",
        ],
        acknowledgments="Supported by NIH Grant R01-CA123456",
        contact_info="jane.smith@university.edu",
    )
